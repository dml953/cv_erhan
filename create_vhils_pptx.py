from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

CONCRETE = RGBColor(0xD4, 0xCF, 0xC4)
DUST = RGBColor(0xA8, 0x9F, 0x91)
SHADOW = RGBColor(0x1A, 0x1A, 0x1A)
RUST = RGBColor(0x8B, 0x45, 0x13)
DARK_BG = RGBColor(0x12, 0x10, 0x0D)
DARKER_BG = RGBColor(0x0A, 0x09, 0x08)
WARM_DARK = RGBColor(0x1E, 0x1B, 0x17)
CRACK = RGBColor(0x3D, 0x34, 0x28)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_name='Courier New',
                font_size=Pt(18), font_color=CONCRETE, bold=False, alignment=PP_ALIGN.LEFT,
                italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    return txBox, tf


def add_paragraph(tf, text, font_name='Courier New', font_size=Pt(14),
                  font_color=DUST, bold=False, alignment=PP_ALIGN.LEFT,
                  space_before=Pt(6), space_after=Pt(6), italic=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_shape_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


# ══════════════════════════════════════════
# SLIDE 1: TITRE
# ══════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide1, DARK_BG)

add_shape_rect(slide1, Inches(0), Inches(0), SLIDE_W, Inches(0.08), RUST)

add_textbox(slide1, Inches(1), Inches(1.8), Inches(11), Inches(2.5),
            'VHILS', font_size=Pt(120), font_color=CONCRETE, bold=True,
            alignment=PP_ALIGN.CENTER)

_, tf = add_textbox(slide1, Inches(2), Inches(4.2), Inches(9), Inches(0.8),
                    '——— L\'ART DE LA DESTRUCTION ———', font_size=Pt(22),
                    font_color=RUST, alignment=PP_ALIGN.CENTER)

_, tf = add_textbox(slide1, Inches(2), Inches(5.3), Inches(9), Inches(0.6),
                    'ALEXANDRE FARTO  •  NÉ EN 1987  •  LISBONNE, PORTUGAL',
                    font_size=Pt(12), font_color=DUST, alignment=PP_ALIGN.CENTER)

add_shape_rect(slide1, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), RUST)

# ══════════════════════════════════════════
# SLIDE 2: BIOGRAPHIE
# ══════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, WARM_DARK)

add_shape_rect(slide2, Inches(0), Inches(0), Inches(0.12), SLIDE_H, RUST)

add_textbox(slide2, Inches(0.8), Inches(0.4), Inches(4), Inches(0.5),
            'L\'ARTISTE', font_size=Pt(10), font_color=RUST,
            alignment=PP_ALIGN.LEFT)

add_textbox(slide2, Inches(0.8), Inches(0.9), Inches(5), Inches(1.2),
            'Alexandre\nManuel\nDias Farto', font_size=Pt(32), font_color=CONCRETE,
            bold=True, alignment=PP_ALIGN.LEFT)

facts = [
    ('NAISSANCE', '1987, Seixal, Portugal'),
    ('PSEUDONYME', 'VHILS'),
    ('FORMATION', 'University of the Arts London'),
    ('DISCIPLINE', 'Street Art / Art Urbain'),
    ('MOUVEMENT', 'Art Contemporain Destructif'),
]

y_pos = 3.0
for label, value in facts:
    add_textbox(slide2, Inches(0.8), Inches(y_pos), Inches(4), Inches(0.25),
                label, font_size=Pt(8), font_color=RUST)
    add_textbox(slide2, Inches(0.8), Inches(y_pos + 0.25), Inches(4), Inches(0.35),
                value, font_size=Pt(13), font_color=CONCRETE)
    y_pos += 0.75

add_shape_rect(slide2, Inches(5.5), Inches(0.5), Pt(1), Inches(6.5), CRACK)

_, tf = add_textbox(slide2, Inches(6.2), Inches(0.8), Inches(6.5), Inches(6),
                    '', font_size=Pt(14), font_color=DUST)
tf.paragraphs[0].text = ''

bio_paragraphs = [
    "Alexandre Farto, connu sous le nom de VHILS, est un artiste visuel portugais qui a révolutionné le street art en développant une technique unique de soustraction plutôt que d'addition.",
    "",
    "Grandissant dans la banlieue industrielle de Lisbonne après la Révolution des Œillets, il est fasciné par les murs décrépits et les couches d'affiches qui racontent l'histoire urbaine.",
    "",
    "Dès 13 ans, il commence le graffiti avant de développer sa propre voie artistique radicalement différente.",
    "",
    "En 2008, une de ses œuvres est exposée à côté d'une pièce de Banksy lors du Cans Festival à Londres, propulsant sa carrière internationale.",
]

for para in bio_paragraphs:
    add_paragraph(tf, para, font_size=Pt(13), font_color=DUST,
                  space_before=Pt(4), space_after=Pt(4))


# ══════════════════════════════════════════
# SLIDE 3: PROCESSUS CRÉATIF
# ══════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3, DARKER_BG)

add_textbox(slide3, Inches(0.8), Inches(0.5), Inches(8), Inches(1),
            'PROCESSUS CRÉATIF', font_size=Pt(40), font_color=CONCRETE,
            bold=True)

add_shape_rect(slide3, Inches(0.8), Inches(1.4), Inches(1.5), Pt(4), RUST)

techniques = [
    ('⚒  GRAVURE MURALE', 'À l\'aide de burins et de marteaux, VHILS creuse directement dans les murs, retirant couche après couche d\'enduit, de plâtre et de peinture pour révéler des portraits monumentaux.'),
    ('💥  EXPLOSIFS', 'Pour certaines œuvres à grande échelle, il utilise des micro-charges explosives contrôlées qui détachent la surface du mur selon un motif prédéfini, créant des portraits en quelques secondes.'),
    ('⚡  ACIDE & EAU DE JAVEL', 'Sur des surfaces métalliques ou des panneaux publicitaires, il applique des agents chimiques qui rongent sélectivement la matière, créant des dégradés subtils et des textures organiques.'),
]

x_positions = [0.5, 4.5, 8.7]
for i, (title, desc) in enumerate(techniques):
    x = x_positions[i]
    card = add_shape_rect(slide3, Inches(x), Inches(2.2), Inches(3.8), Inches(4.8),
                          RGBColor(0x1E, 0x1B, 0x17), CRACK)

    add_textbox(slide3, Inches(x + 0.3), Inches(2.5), Inches(3.2), Inches(0.5),
                f'0{i+1}', font_size=Pt(48), font_color=RGBColor(0x2A, 0x25, 0x20),
                bold=True)

    add_textbox(slide3, Inches(x + 0.3), Inches(3.3), Inches(3.2), Inches(0.6),
                title, font_size=Pt(11), font_color=RUST, bold=True)

    add_textbox(slide3, Inches(x + 0.3), Inches(4.0), Inches(3.2), Inches(2.8),
                desc, font_size=Pt(11), font_color=DUST)


# ══════════════════════════════════════════
# SLIDE 4: ŒUVRES MAJEURES
# ══════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4, RGBColor(0x14, 0x12, 0x0F))

add_textbox(slide4, Inches(0.8), Inches(0.4), Inches(8), Inches(1),
            'ŒUVRES MAJEURES', font_size=Pt(38), font_color=CONCRETE, bold=True)
add_shape_rect(slide4, Inches(0.8), Inches(1.25), Inches(1.5), Pt(4), RUST)

works = [
    ('2008', 'SCRATCHING THE SURFACE',
     'Série fondatrice où VHILS grave des portraits dans les murs de Lisbonne, révélant les strates historiques de la ville.',
     'Lisbonne, Portugal'),
    ('2014', 'DISSECTION',
     'Installation immersive explorant la déconstruction de l\'identité à travers des portraits sculptés en couches de matériaux urbains.',
     'Hong Kong'),
    ('2019', 'ETHEREAL',
     'Série utilisant la pyrotechnie et les explosifs pour créer des portraits éphémères capturés en vidéo haute vitesse.',
     'International'),
    ('2021', 'HOMENAGEM AOS PROFISSIONAIS DE SAÚDE',
     'Immense portrait d\'une soignante gravé sur un mur d\'hôpital en hommage aux travailleurs de santé pendant la pandémie.',
     'Hospital São José, Lisbonne'),
]

y_positions = [1.8, 3.2, 4.6, 6.0]
for i, (year, title, desc, location) in enumerate(works):
    y = y_positions[i]

    add_shape_rect(slide4, Inches(0.8), Inches(y), Pt(3), Inches(1.2), RUST)

    add_textbox(slide4, Inches(1.1), Inches(y), Inches(1.2), Inches(0.4),
                year, font_size=Pt(10), font_color=RUST)

    add_textbox(slide4, Inches(2.5), Inches(y - 0.05), Inches(5), Inches(0.45),
                title, font_size=Pt(14), font_color=CONCRETE, bold=True)

    add_textbox(slide4, Inches(2.5), Inches(y + 0.4), Inches(6), Inches(0.7),
                desc, font_size=Pt(11), font_color=DUST)

    add_textbox(slide4, Inches(9.5), Inches(y + 0.15), Inches(3.5), Inches(0.4),
                f'📍 {location}', font_size=Pt(9), font_color=CRACK)


# ══════════════════════════════════════════
# SLIDE 5: PHILOSOPHIE / CITATION
# ══════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5, RGBColor(0x15, 0x12, 0x0E))

add_textbox(slide5, Inches(0.5), Inches(0.8), Inches(2), Inches(2),
            '«', font_size=Pt(200), font_color=RGBColor(0x3A, 0x20, 0x10),
            bold=True)

add_textbox(slide5, Inches(1.5), Inches(2.0), Inches(10), Inches(2.5),
            'Je ne crée pas en ajoutant,\nje crée en enlevant.\nChaque coup de burin révèle\nce qui était déjà là,\ncaché sous la surface.',
            font_size=Pt(28), font_color=CONCRETE, italic=True,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide5, Inches(2), Inches(4.8), Inches(9), Inches(0.5),
            '— VHILS', font_size=Pt(14), font_color=RUST,
            alignment=PP_ALIGN.CENTER)

add_shape_rect(slide5, Inches(4), Inches(5.5), Inches(5), Pt(1), CRACK)

add_textbox(slide5, Inches(2.5), Inches(5.8), Inches(8), Inches(1.5),
            'Sa philosophie artistique repose sur l\'idée que la beauté existe déjà dans les murs de nos villes. L\'artiste ne fait que la libérer en retirant les couches superflues, comme un archéologue urbain mettant au jour les mémoires enfouies de l\'espace public.',
            font_size=Pt(13), font_color=DUST, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════
# SLIDE 6: RAYONNEMENT INTERNATIONAL
# ══════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide6, RGBColor(0x16, 0x13, 0x10))

add_textbox(slide6, Inches(0.8), Inches(0.4), Inches(10), Inches(1),
            'RAYONNEMENT MONDIAL', font_size=Pt(38), font_color=CONCRETE, bold=True)
add_shape_rect(slide6, Inches(0.8), Inches(1.25), Inches(1.5), Pt(4), RUST)

cities = [
    ('LISBONNE', 'Portugal'),
    ('LONDRES', 'Royaume-Uni'),
    ('PARIS', 'France'),
    ('RIO', 'Brésil'),
    ('HONG KONG', 'Chine'),
    ('LOS ANGELES', 'États-Unis'),
    ('SYDNEY', 'Australie'),
    ('MOSCOU', 'Russie'),
]

cols = 4
rows = 2
cell_w = 2.8
cell_h = 2.2
start_x = 0.8
start_y = 2.0
gap_x = 0.3
gap_y = 0.3

for i, (city, country) in enumerate(cities):
    col = i % cols
    row = i // cols
    x = start_x + col * (cell_w + gap_x)
    y = start_y + row * (cell_h + gap_y)

    add_shape_rect(slide6, Inches(x), Inches(y), Inches(cell_w), Inches(cell_h),
                   RGBColor(0x1E, 0x1B, 0x17), CRACK)

    add_textbox(slide6, Inches(x), Inches(y + 0.5), Inches(cell_w), Inches(0.5),
                '◉', font_size=Pt(24), font_color=RUST, alignment=PP_ALIGN.CENTER)

    add_textbox(slide6, Inches(x), Inches(y + 1.1), Inches(cell_w), Inches(0.5),
                city, font_size=Pt(13), font_color=CONCRETE, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide6, Inches(x), Inches(y + 1.5), Inches(cell_w), Inches(0.4),
                country, font_size=Pt(9), font_color=DUST,
                alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════
# SLIDE 7: EXPOSITIONS CLÉS
# ══════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide7, RGBColor(0x11, 0x0F, 0x0C))

add_textbox(slide7, Inches(0.8), Inches(0.4), Inches(8), Inches(1),
            'EXPOSITIONS CLÉS', font_size=Pt(38), font_color=CONCRETE, bold=True)
add_shape_rect(slide7, Inches(0.8), Inches(1.25), Inches(1.5), Pt(4), RUST)

add_shape_rect(slide7, Inches(1.5), Inches(1.8), Pt(2), Inches(5.3), RUST)

expos = [
    ('2008', 'Cans Festival', 'Leake Street Tunnel, Londres — aux côtés de Banksy'),
    ('2011', 'Husk', 'Lazarides Gallery, Londres — première expo solo majeure'),
    ('2014', 'Vestiges', 'Musée EDP, Lisbonne — rétrospective institutionnelle'),
    ('2016', 'Annihilation', 'Palais de Tokyo, Paris'),
    ('2019', 'Metamorphosis', 'MAAT, Lisbonne — exploration multimédia'),
    ('2022', 'Imprint', 'Hong Kong Contemporary Art Foundation'),
]

y_pos = 2.0
for year, name, place in expos:
    diamond = slide7.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(1.35), Inches(y_pos + 0.1),
                                       Pt(12), Pt(12))
    diamond.fill.solid()
    diamond.fill.fore_color.rgb = RUST
    diamond.line.fill.background()

    add_textbox(slide7, Inches(2.0), Inches(y_pos - 0.05), Inches(1.5), Inches(0.4),
                year, font_size=Pt(10), font_color=RUST)

    add_textbox(slide7, Inches(3.3), Inches(y_pos - 0.08), Inches(5), Inches(0.45),
                name, font_size=Pt(16), font_color=CONCRETE, bold=True)

    add_textbox(slide7, Inches(3.3), Inches(y_pos + 0.35), Inches(8), Inches(0.4),
                place, font_size=Pt(11), font_color=DUST)

    y_pos += 0.9


# ══════════════════════════════════════════
# SLIDE 8: CONCLUSION
# ══════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide8, RGBColor(0x14, 0x11, 0x0D))

add_shape_rect(slide8, Inches(0), Inches(0), SLIDE_W, Inches(0.08), RUST)

add_textbox(slide8, Inches(1.5), Inches(1.0), Inches(10), Inches(0.6),
            'EN RÉSUMÉ', font_size=Pt(12), font_color=RUST,
            alignment=PP_ALIGN.CENTER)

add_textbox(slide8, Inches(1.5), Inches(1.6), Inches(10), Inches(1.5),
            'DÉTRUIRE\nPOUR RÉVÉLER', font_size=Pt(52), font_color=CONCRETE,
            bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide8, Inches(2), Inches(3.8), Inches(9), Inches(2),
            'VHILS redéfinit les frontières du street art en transformant la destruction en acte créatif. Ses œuvres sont des palimpsestes urbains qui questionnent notre rapport à la ville, à la mémoire collective et à l\'identité.\n\nChaque mur qu\'il touche devient un témoignage vivant de notre époque.',
            font_size=Pt(14), font_color=DUST, alignment=PP_ALIGN.CENTER)

links_text = 'vfrfranco.com   •   @vfrfranco   •   Underdogs Gallery'
add_textbox(slide8, Inches(2), Inches(6.3), Inches(9), Inches(0.5),
            links_text, font_size=Pt(11), font_color=CRACK,
            alignment=PP_ALIGN.CENTER)

add_shape_rect(slide8, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), RUST)


# SAVE
prs.save('/workspace/VHILS_Presentation.pptx')
print("✓ Fichier VHILS_Presentation.pptx créé avec succès")
